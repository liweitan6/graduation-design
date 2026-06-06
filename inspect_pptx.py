import os
from pptx import Presentation

def extract_pptx_outline(filepath, output_txt_path):
    if not os.path.exists(filepath):
        print(f"Error: {filepath} does not exist.")
        return
    prs = Presentation(filepath)
    with open(output_txt_path, 'w', encoding='utf-8') as f:
        f.write(f"=== Outline for {filepath} ===\n\n")
        for i, slide in enumerate(prs.slides):
            f.write(f"--- Slide {i+1} ---\n")
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            f.write(f"Title: {title}\n")
            f.write("Text content:\n")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if slide.shapes.title and shape == slide.shapes.title:
                        continue
                    text = shape.text_frame.text.strip()
                    if text:
                        # Indent lines for readability
                        indented_lines = "\n".join(["  " + line for line in text.split("\n") if line.strip()])
                        f.write(indented_lines + "\n")
            f.write("\n")

extract_pptx_outline("毕业设计答辩-参考模板（王冠琦-20373008-排球比赛视频行为识别算法研究）第二三部分可按两或三个研究内容介绍（每个含研究工作和对应实验结果）.pptx", "ref_template_outline.txt")
extract_pptx_outline("22371240-郭俊杰-答辩.pptx", "user_draft_outline.txt")
print("Done extracting outlines.")
